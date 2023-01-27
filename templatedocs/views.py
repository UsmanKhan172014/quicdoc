from django.shortcuts import render
from django.http import HttpResponse, HttpResponseRedirect, FileResponse
# from .forms import UploadFileForm
from .models import UploadFileForm
import docx
import os
import re
from pptx import Presentation
from django.contrib.auth.decorators import login_required
from tika import parser


def file_to_placeholders_array(file_path):
    text = parser.from_file(file_path)
    text = re.sub(r'\s+', ' ', text['content'])
    placeholders = re.findall(r"\{[^}]*\}", text)
    placeholders = list(dict.fromkeys(placeholders))
    print(placeholders)
    return placeholders


@login_required
def home(request):
    current_user = request.user
    if request.method == 'POST':
        form = UploadFileForm(request.POST, request.FILES)
        file = request.FILES['file']
        data = UploadFileForm.objects.create(name=request.POST['name'],
                                             description=request.POST['description'],
                                             file=file,
                                             user_id=current_user.id
                                             )
        data.save()
        return HttpResponseRedirect('/templatedocs/get_all/')

    else:
        form = UploadFileForm()
    return render(request, 'templatedocs/upload.html', {'form': form})


# home index function
@login_required
def get_all(request):
    current_user = request.user
    data = UploadFileForm.objects.filter(user_id=current_user.id)
    return render(request, 'templatedocs/index.html', {
        'data': data,
        'active_tab': 'get_all',
    })


# delete function
@login_required
def delete(request, id):
    data = UploadFileForm.objects.get(id=id)
    # find the file path
    file_path = data.file.path
    # delete the file
    os.remove(file_path)
    # delete the data from database
    data.delete()
    return HttpResponseRedirect('/templatedocs/get_all/')


@login_required
def updateDoc(request, id):
    data = UploadFileForm.objects.get(id=id)
    # baseurl = request.build_absolute_uri('/')
    base_path = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    file = base_path + "/media/" + str(data.file)
    place_holders = file_to_placeholders_array(file)
    return render(request, 'templatedocs/updateDoc.html', {'id': id, 'unique_placeholders': place_holders})
    # if ".docx" in file:
    #     text_file_path = extract_word_to_text(file)
    #     print(text_file_path)
    #     placeholders = get_text_placeholders(text_file_path)
    #     # print(placeholders)
    #     return render(request, 'templatedocs/updateDoc.html', {'id': id, 'unique_placeholders': placeholders})
    #
    #     doc = docx.Document(data.file)
    #     unique_placeholders = set()
    #     for para in doc.paragraphs:
    #         placeholders = re.findall(r"\{[^}]*\}", para.text)
    #         unique_placeholders.update(placeholders)
    #
    #     for table in doc.tables:
    #         for row in table.rows:
    #             for cell in row.cells:
    #                 for para in cell.paragraphs:
    #                     # need code to find multiple items in a single string
    #                     placeholders = re.findall(r"\{[^}]*\}", para.text)
    #                     unique_placeholders.update(placeholders)
    #
    #     # print(unique_placeholders)
    #     unique_placeholders = list(unique_placeholders)
    #     # print(unique_placeholders)
    #     return render(request, 'templatedocs/updateDoc.html', {'id': id, 'unique_placeholders': unique_placeholders})
    # elif ".pptx" in file:
    #     prs = Presentation(file)
    #     unique_placeholders = set()
    #     for slide in prs.slides:
    #         for shape in slide.shapes:
    #             if not shape.has_text_frame:
    #                 continue
    #             for paragraph in shape.text_frame.paragraphs:
    #                 placeholders = re.findall(r"\{[^}]*\}", paragraph.text)
    #                 print(placeholders)
    #                 unique_placeholders.update(placeholders)
    #     print(unique_placeholders)
    #     unique_placeholders = list(unique_placeholders)
    #     print(unique_placeholders)
    #     # return HttpResponse(unique_placeholders)
    #     return render(request, 'templatedocs/updateDoc.html', {'id': id, 'unique_placeholders': unique_placeholders})
    # else:
    #     return HttpResponse("File not supported")


@login_required
def update(request, id):
    data = UploadFileForm.objects.get(id=id)
    if request.method == 'POST':
        base_path = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        if ".docx" in data.file.name:
            doc = docx.Document(data.file)

            try:
                line_count = 1
                for para in doc.paragraphs:
                    placeholders = re.findall(r"\{[^}]*\}", para.text)
                    for placeholder in placeholders:
                        style1 = doc.styles['heading 1']
                        style2 = doc.styles['heading 2']
                        default_color = style2.font.color.rgb
                        if str(default_color) == 'FFFFFF':
                            para.text = para.text.replace(placeholder, request.POST[placeholder])
                            para.style.font.color.rgb = style1.font.color.rgb
                        else:
                            para.text = para.text.replace(placeholder, request.POST[placeholder])
                        line_count += 1
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            for para in cell.paragraphs:
                                placeholders = re.findall(r"\{[^}]*\}", para.text)
                                for placeholder in placeholders:
                                    style = para.style
                                    para.text = para.text.replace(placeholder, request.POST[placeholder])
                                    para.style = style
                new_name = str(data.file).replace(".docx", "_updated.docx")
                doc.save(base_path + "/media/documents/updated.docx")
                response = FileResponse(open(base_path + "/media/documents/updated.docx", 'rb'),
                                        content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
                response['Content-Disposition'] = 'attachment; filename="{}"'.format(new_name)
                return response
            except Exception as e:
                print(e)
                return HttpResponse(e)
        elif ".pptx" in data.file.name:
            prs = Presentation(base_path + "/media/" + str(data.file))
            try:
                slide_number = 1
                for slide in prs.slides:
                    shape_number = 1
                    for shape in slide.shapes:
                        text_frames = 1
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            placeholders = re.findall(r"\{[^}]*\}", paragraph.text)
                            for placeholder in placeholders:
                                for run in paragraph.runs:
                                    font_name = run.font.name
                                    font_size = run.font.size
                                    paragraph.text = paragraph.text.replace(placeholder, request.POST[placeholder])
                                    paragraph.runs[0].font.name = font_name
                                    paragraph.runs[0].font.size = font_size
                                    print("Slide: " + str(slide_number) + " Shape: " + str(
                                        shape_number) + " Text Frame: " + str(text_frames) + " Paragraph: " + str(
                                        paragraph.text))
                            text_frames += 1
                        shape_number += 1
                    slide_number += 1

                new_name = str(data.file).replace(".pptx", "_updated.pptx")
                prs.save(base_path + "/media/documents/updated.pptx")
                response = FileResponse(open(base_path + "/media/documents/updated.pptx", 'rb'),
                                        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
                response['Content-Disposition'] = 'attachment; filename="{}"'.format(new_name)
                return response
                # return HttpResponse("Success")
            except Exception as e:
                print(e)
                return HttpResponse("Error")
        else:
            return HttpResponse("File not supported")
    else:
        return HttpResponseRedirect('/templatedocs/get_all/')
